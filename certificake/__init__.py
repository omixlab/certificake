from pptx import Presentation
from jinja2 import Environment, BaseLoader
import sys
import argparse
import subprocess
import pandas as pd
import os

def generate_certificate(template_file, data, output_base, pdf=False):

    data = data.to_dict()

    prs = Presentation(template_file)
    
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        text_template = Environment(loader=BaseLoader).from_string(run.text)
                        run.text = text_template.render(**data)
    
    output_pptx_raw = output_base + '.pptx'

    output_pptx_formater = Environment(loader=BaseLoader).from_string(output_pptx_raw)
    output_pptx_repr = output_pptx_formater.render(**data)
    output_pptx_repr_dirname = os.path.dirname(output_pptx_repr)
    os.makedirs(output_pptx_repr_dirname, exist_ok=True)

    prs.save(output_pptx_repr)

    if pdf:
        generate_pdf(output_pptx_repr)

def generate_pdf(pptx_file):
    current_dir = os.getcwd()
    output_dir = os.path.dirname(pptx_file)
    output_pptx_repr_basename = os.path.basename(pptx_file)
    os.chdir(output_dir)
    if sys.platform == "darwin":
        subprocess.call(f'/Applications/LibreOffice.app/Contents/MacOS/soffice --headless --invisible --convert-to pdf {output_pptx_repr_basename}', shell=True)
    else:
        subprocess.call(f'libreoffice --headless --invisible --convert-to pdf {output_pptx_repr_basename}', shell=True)
    os.chdir(current_dir)


def main():

    arguments_parser = argparse.ArgumentParser()
    arguments_parser.add_argument('--template', help="A PPTX file containing the certificate template", required=True)
    arguments_parser.add_argument('--data', help='A CSV/XLSX file containing the participants data', required=True)
    arguments_parser.add_argument('--output', required="The output directory (will be created with doesn't exist)", required=True)
    arguments_parser.add_argument('--pdf', action='store_true', help='Save as PDF (default is PPTX)')

    arguments = arguments_parser.parse_args()

    _, data_file_ext = os.path.splitext(arguments.data)

    if data_file_ext.lower() == '.csv':
        df_data = pd.read_csv(arguments.data)
    
    elif data_file_ext.lower() == 'xlsx':
        df_data = pd.read_excel(arguments.data)
    
    else:
        print("Error: --data file must be .csv or .xlsx")
        exit(1)

    for r, row in df_data.iterrows():
        generate_certificate(arguments.template, row, arguments.output, pdf=arguments.pdf)

if __name__ == '__main__':
    main()
